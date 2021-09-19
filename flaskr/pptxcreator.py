from flask import url_for
from pptx import Presentation


THEME_PATH = {
    1: 'static/themes/gallery.pptx',
    2: 'static/themes/ion.pptx',
    3: 'static/themes/ion-boardroom.pptx',
    4: 'static/themes/organic.pptx',
    5: 'static/themes/slice.pptx'
}


def create_slides(theme, name, titles, summaries):

    """
    Inputs:
        name -- title slide name
        titles -- (str) slide titles
        summaries -- (list[list]) slide content
    """

    prs = Presentation(THEME_PATH[theme])

    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = name
    subtitle = slide.placeholders[1]
    subtitle.text = "Made by EZ PowerPoint"

    for i in range(len(titles)):

        slide = prs.slides.add_slide(bullet_slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = titles[i]

        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame

        for index, point in enumerate(summaries[i]):
            if index == 0:
                tf.text = point
            else:
                p = tf.add_paragraph()
                p.text = point

    prs.save(name + '.pptx')
